# -*- coding: utf-8 -*-
"""
make_sample_data.py
===================
코드가 잘 도는지 확인하기 위한 가상 데이터를 data/ 폴더에 만든다.
실제 데이터가 준비되면 필요 없다(그때는 data/ 안의 sample_* 파일을 지우면 된다).

만드는 것:
  sample_xlsx_gt.xlsx / sample_xlsx_before.txt / sample_xlsx_after.txt
     - 엑셀 GT + 파서 결과(마크다운 파이프 표, null 패딩·Unnamed 흔적 포함)
       before: 헤더를 Unnamed로 날리고 값 하나 오인식, 패딩 행 잔뜩
       after : 값·헤더 정상
  sample_ppt_gt.pptx / sample_ppt_before.txt / sample_ppt_after.txt
     - 3장짜리 PPT + 파서 결과(## 로 슬라이드 시작 표시)
       before: 2번 슬라이드가 통째로 누락(뒤가 밀리는 상황 재현)
       after : 3장 모두 정상
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook

import config

DATA = config.DATA_DIR
DATA.mkdir(parents=True, exist_ok=True)


def _make_xlsx():
    wb = Workbook(); ws = wb.active; ws.title = "실적"
    ws.append(["구분", "202603", "202604"])
    ws.append(["저축은행", 429, 322])
    ws.append(["인터넷은행", 3329, 2605])
    wb.save(DATA / "sample_xlsx_gt.xlsx")

    # 파서 before: 헤더 Unnamed, 값 오인식(322→320), null 패딩 다수
    (DATA / "sample_xlsx_before.txt").write_text(
        "## 헤더구조,음영추출\n\n"
        "| 구분 | Unnamed: 1 | Unnamed: 2 |\n"
        "| --- | --- | --- |\n"
        "| 저축은행 | 429 | 320 |\n"
        "| 인터넷은행 | 3329 | 2605 |\n"
        "| null | null | null |\n"
        "| null | null | null |\n"
        "| null | null | null |\n",
        encoding="utf-8")

    # 파서 after: 헤더·값 정상
    (DATA / "sample_xlsx_after.txt").write_text(
        "| 구분 | 202603 | 202604 |\n"
        "| 저축은행 | 429 | 322 |\n"
        "| 인터넷은행 | 3329 | 2605 |\n"
        "| null | null |\n",
        encoding="utf-8")


def _add_slide(prs, title, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4)).text_frame
    tb.text = body_lines[0]
    for line in body_lines[1:]:
        p = tb.add_paragraph(); p.text = line; p.level = 1


def _make_pptx():
    prs = Presentation()
    _add_slide(prs, "JURIS 준법심의 매뉴얼", ["준법심의 대상은 신상품 및 광고물이다", "심의는 출시 전에 완료한다"])
    _add_slide(prs, "특별점검 절차", ["정기점검과 특별점검으로 나뉜다", "현장점검은 분기 1회 실시한다"])
    _add_slide(prs, "유의 모집인 관리", ["유의 모집인은 별도 관리한다", "위반 시 자격을 정지한다"])
    prs.save(DATA / "sample_ppt_gt.pptx")

    # 파서 before: 2번 슬라이드(특별점검)가 통째로 누락 → 통짜로 이으면 뒤가 밀림
    (DATA / "sample_ppt_before.txt").write_text(
        "## JURIS 준법심의 매뉴얼\n"
        "준법심의 대상은 신상품 및 광고물이다\n"
        "- 심의는 출시 전에 완료한다\n\n"
        "## 유의 모집인 관리\n"
        "유의 모집인은 별도 관리한다\n"
        "- 위반 시 자격을 정지한다\n",
        encoding="utf-8")

    # 파서 after: 3장 모두 정상
    (DATA / "sample_ppt_after.txt").write_text(
        "## JURIS 준법심의 매뉴얼\n"
        "준법심의 대상은 신상품 및 광고물이다\n"
        "- 심의는 출시 전에 완료한다\n\n"
        "## 특별점검 절차\n"
        "정기점검과 특별점검으로 나뉜다\n"
        "- 현장점검은 분기 1회 실시한다\n\n"
        "## 유의 모집인 관리\n"
        "유의 모집인은 별도 관리한다\n"
        "- 위반 시 자격을 정지한다\n",
        encoding="utf-8")


def make():
    _make_xlsx()
    _make_pptx()
    print(f"[완료] 가상 데이터 생성 → {DATA}/  (sample_xlsx_*, sample_ppt_*)")


if __name__ == "__main__":
    make()