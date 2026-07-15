# -*- coding: utf-8 -*-
"""
gt_extract.py
=============
GT(정답) 문서에서 "본문 텍스트"와 "표 목록"을 뽑아내는 통합 추출기.

GT는 .docx / .pptx / .xlsx 세 가지 형식으로 올 수 있고, 확장자를 보고 알맞은
추출기로 자동 분기한다. 어떤 형식이든 결과는 항상 아래 두 가지다.

  - text  : 표를 뺀 본문 텍스트(원문 순서대로 이어붙임).
            제목/볼드/취소선/리스트 같은 서식은 파서 출력과 같은 Markdown 기호
            (## , **..** , ~~..~~ , - )로 "복원"해서 넣는다. → WER 비교용
  - tables: 문서에 나온 표들을 등장 순서대로 <table> HTML 문자열로 변환한 리스트.
            → TEDS-content 비교용

왜 서식을 Markdown 기호로 복원하는가
------------------------------------
파서 출력(txt)은 제목을 "# ", 볼드를 "**..**", 취소선을 "~~..~~", 목록을 "- "처럼
Markdown 기호로 표현한다. GT 원본에는 이런 기호가 없고 Word/PPT의 서식 속성(굵게,
취소선, 제목 스타일, 목록 스타일)으로만 들어있다. GT를 순수 텍스트로만 뽑으면 파서가
서식을 정확히 살려도 그 기호가 전부 삽입 오류로 잡힌다. 그래서 GT도 파서와 "같은
기호 체계"로 맞춰줘야 공정하게 비교된다(normalize.py의 공정 비교 원칙 참고).

서식 → 기호 변환 규칙
---------------------
  - 제목 스타일("Heading N" / "제목 N")   → "#" * N + 텍스트
  - run 단위 굵게(bold)                    → **텍스트**
  - run 단위 취소선(strike)                → ~~텍스트~~   (bold와 겹치면 **~~텍스트~~**)
  - 목록 스타일 문단(List/Bullet/번호)     → "- " + 텍스트
  (표 안 셀 텍스트는 단순화를 위해 서식 복원 없이 순수 텍스트만 사용한다.
   표는 TEDS-content가 별도로 구조+셀 내용을 비교한다.)

한계
----
  - 표의 병합 셀(rowspan/colspan)은 복원하지 않는다(셀 텍스트 + 기본 행/열 구조까지).
  - 도형/텍스트 상자/SmartArt 안의 텍스트는 추출하지 않는다(문단·표만 순회).
  - 순서 있는 목록의 실제 번호는 복원하지 않고 "- "로 통일한다.
"""

import re
import html as html_lib
from pathlib import Path


# ══════════════════════════════════════════════════════════════════════
# 공통 유틸: run(서식 최소 단위) → Markdown 기호로 복원
# ══════════════════════════════════════════════════════════════════════
def _wrap_run_markdown(text: str, bold: bool, strike: bool) -> str:
    """한 run의 텍스트를 bold/strike 여부에 따라 Markdown 기호로 감싼다."""
    if not text:
        return ""
    if strike:
        text = f"~~{text}~~"
    if bold:
        text = f"**{text}**"
    return text


def _run_bold_strike(run):
    """
    run 하나의 (bold, strike) 여부를 docx/pptx 양쪽에서 안전하게 읽는다.
      - bold  : docx는 run.bold, pptx는 run.font.bold
      - strike: docx는 run.font.strike, pptx는 rPr XML의 strike 속성(noStrike가 아니면 True)
    """
    font = getattr(run, "font", None)

    bold = getattr(run, "bold", None)                 # docx Run
    if bold is None and font is not None:
        bold = getattr(font, "bold", None)            # pptx Run.font

    strike = getattr(font, "strike", None) if font is not None else None  # docx Font.strike
    if strike is None and font is not None:
        rpr = getattr(font, "_rPr", None)             # pptx: 문자 속성 XML
        if rpr is not None:
            val = rpr.get("strike")
            strike = bool(val) and val != "noStrike"

    return bool(bold), bool(strike)


def _runs_to_markdown(runs) -> str:
    """
    run 목록을 받아 bold/strike가 같은 연속 run은 하나로 합친 뒤 기호로 감싼다.
    (연속된 같은 서식을 합쳐 "**a****b**"처럼 기호가 겹치는 것을 막는다.)
    """
    groups = []  # [(합친 텍스트, bold, strike), ...]
    for run in runs:
        text = run.text or ""
        if not text:
            continue
        bold, strike = _run_bold_strike(run)
        if groups and groups[-1][1] == bold and groups[-1][2] == strike:
            g = groups[-1]
            groups[-1] = (g[0] + text, bold, strike)
        else:
            groups.append((text, bold, strike))
    return "".join(_wrap_run_markdown(t, b, s) for t, b, s in groups)


# ══════════════════════════════════════════════════════════════════════
# .docx 추출
# ══════════════════════════════════════════════════════════════════════
from docx import Document as _DocxDocument
from docx.table import Table as _DocxTable
from docx.text.paragraph import Paragraph as _DocxParagraph
from docx.oxml.text.paragraph import CT_P as _CT_P
from docx.oxml.table import CT_Tbl as _CT_Tbl

# "Heading 2" / "제목 2" → 레벨 2
_HEADING_STYLE = re.compile(r"(?:Heading|제목)\s*(\d+)", re.IGNORECASE)
# 목록 계열 스타일 이름(영문/한글판)
_LIST_STYLE = re.compile(r"(List|Bullet|Number|목록|글머리)", re.IGNORECASE)


def _docx_iter_blocks(document):
    """문서 본문을 순서대로 돌며 문단/표를 구분해 내보낸다(python-docx 미제공 기능)."""
    for child in document.element.body.iterchildren():
        if isinstance(child, _CT_P):
            yield _DocxParagraph(child, document)
        elif isinstance(child, _CT_Tbl):
            yield _DocxTable(child, document)


def _docx_paragraph_to_markdown(paragraph) -> str:
    """docx 문단 하나를 서식(제목/목록/볼드/취소선) 복원된 텍스트로 변환."""
    if not paragraph.text.strip():
        return ""

    style_name = (paragraph.style.name if paragraph.style else "") or ""

    # 제목 스타일이면 "#" 레벨로
    m = _HEADING_STYLE.match(style_name)
    if m:
        level = min(int(m.group(1)), 6)
        return f"{'#' * level} {paragraph.text.strip()}"

    # run 단위로 볼드/취소선 복원
    body = _runs_to_markdown(paragraph.runs).strip()

    # 목록 스타일이면 "- " 접두어
    if _LIST_STYLE.search(style_name):
        return f"- {body}"
    return body


def _docx_table_to_html(table) -> str:
    """docx 표 → <table> HTML(셀 텍스트만, 병합 미복원)."""
    rows_html = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            txt = "\n".join(p.text for p in cell.paragraphs).strip()
            cells.append(f"<td>{html_lib.escape(txt)}</td>")
        rows_html.append(f"<tr>{''.join(cells)}</tr>")
    return f"<table>{''.join(rows_html)}</table>"


def _docx_textbox_texts(document) -> list:
    """
    본문 문단/표 순회로는 안 잡히는 텍스트박스(VML v:textbox 또는 DrawingML wps:txbx)
    안의 텍스트를 모두 뽑는다.

    왜 필요한가: "Word 텍스트박스 사용 지양" 가이드라인(before=텍스트박스 사용,
    after=표/텍스트로 통합)을 채점하려면, GT에는 "그 정보가 존재한다"는 사실 자체가
    있어야 한다. 이전 구현은 문단/표만 순회해 텍스트박스 내용이 GT에서 통째로
    빠졌었다 - before 문서를 채점해도 파서가 그 내용을 놓쳤는지 알 방법이 없었다.
    """
    body = document.element.body
    out = []
    for el in body.iter():
        if el.tag.endswith("}txbxContent"):
            texts = [t.text for t in el.iter() if t.tag.endswith("}t") and t.text]
            joined = "".join(texts).strip()
            if joined:
                out.append(joined)
    return out


def _docx_comment_texts(document) -> list:
    """
    문서에 달린 메모(comment) 텍스트를 모두 뽑는다.
    "메모 사용 지양(핵심 정보는 본문에 기재)" 가이드라인 채점용 - 위 텍스트박스와
    같은 이유로, GT가 메모 내용의 존재를 모르면 before 문서에서 파서가 메모를
    놓친 것을 점수에 반영할 수 없다. python-docx>=1.1의 Document.comments API 사용.
    """
    comments = getattr(document, "comments", None)
    if not comments:
        return []
    out = []
    for c in comments:
        t = (c.text or "").strip()
        if t:
            out.append(t)
    return out


def extract_docx(path: str, include_textboxes: bool = True, include_comments: bool = True):
    """
    docx GT 추출. 반환은 여전히 (text, tables)지만, text 끝에 텍스트박스/메모 내용이
    있으면 "[텍스트박스] ..." / "[메모] ..." 형태로 라벨을 붙여 덧붙인다.
    (본문 순서 안에 정확히 끼워 넣지 않는 이유: 텍스트박스/메모는 본문 흐름과
    무관하게 앵커링되어 원래도 "논리적 위치"가 모호하다 - 이는 그 자체가
    "본문-참조 분리 배치" 가이드라인이 지적하는 문제이기도 하다. 여기서는 "내용의
    존재 여부"를 채점 대상에 포함시키는 것이 목적이므로, 문서 끝에 모아 붙인다.)
    """
    document = _DocxDocument(path)
    texts, tables = [], []
    for block in _docx_iter_blocks(document):
        if isinstance(block, _DocxParagraph):
            md = _docx_paragraph_to_markdown(block)
            if md:
                texts.append(md)
        else:
            tables.append(_docx_table_to_html(block))

    if include_textboxes:
        for tb in _docx_textbox_texts(document):
            texts.append(f"[텍스트박스] {tb}")
    if include_comments:
        for cm in _docx_comment_texts(document):
            texts.append(f"[메모] {cm}")

    return "\n".join(texts), tables


# ══════════════════════════════════════════════════════════════════════
# .pptx 추출
# ══════════════════════════════════════════════════════════════════════
def _para_font_pt(para):
    """
    한 문단의 대표 폰트 크기(pt)를 최선의 방법으로 추정한다.
      1) run들에 명시된 크기가 있으면 그 중 최댓값
      2) 없으면 문단 레벨(para.font.size)
      3) 그래도 없으면 None (크기 정보 없음)
    (레이아웃/마스터 상속 크기는 python-pptx가 직접 안 풀어줘서 여기선 None 처리.)
    """
    sizes = []
    for r in para.runs:
        sz = getattr(r.font, "size", None)
        if sz is not None:
            try:
                sizes.append(sz.pt)
            except AttributeError:
                pass
    if sizes:
        return max(sizes)
    psz = getattr(para.font, "size", None)
    if psz is not None:
        try:
            return psz.pt
        except AttributeError:
            return None
    return None


def extract_pptx_slides(path: str):
    """
    pptx를 '슬라이드별 텍스트 리스트'와 '전체 표 리스트'로 반환한다.
    슬라이드 정렬 비교(slide_eval)에서 슬라이드 경계를 알기 위해 쓴다.

    ── 제목 판단 규칙(파서와 맞춤) ──────────────────────────────────────
    파서는 "슬라이드에서 글자 크기가 가장 큰 텍스트"에 "## "를 붙인다. 그래서 GT도
    슬라이드 안 모든 문단의 폰트 크기를 보고, 최댓값을 가진 문단(들)을 제목으로 보고
    "## "를 붙인다.
      - 제목 문단에 볼드가 있어도 "## "만 붙이고 볼드(**)는 생략한다(파서 규칙).
      - 제목이 아닌 일반 문단의 볼드는 "**..**"로 유지한다.
    폰트 크기가 슬라이드 전체에서 하나도 명시돼 있지 않으면(전부 상속이라 None),
    크기로 못 고르므로 title placeholder를 제목으로 쓰고(있으면), 그것도 없으면
    제목 없이 본문만 낸다.
    """
    from pptx import Presentation

    prs = Presentation(path)
    slides, tables = [], []
    for slide in prs.slides:
        # title placeholder(폴백용)
        title_shape = None
        try:
            if slide.shapes.title is not None:
                title_shape = slide.shapes.title
        except (KeyError, ValueError):
            title_shape = None

        # 1) 슬라이드 안 모든 (문단, shape) 수집 + 폰트 크기 기록
        paras = []  # [(shape, para, text, font_pt)]
        for shape in slide.shapes:
            if shape.has_table:
                tables.append(_pptx_table_to_html(shape.table))
                continue
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                raw = "".join(r.text or "" for r in para.runs)
                if not raw.strip():
                    continue
                paras.append((shape, para, raw, _para_font_pt(para)))

        # 2) 제목 문단 결정
        known_sizes = [fp for *_, fp in paras if fp is not None]
        title_ids = set()   # 제목으로 볼 문단의 id
        if known_sizes:
            max_pt = max(known_sizes)
            # 최대 크기(부동소수 오차 여유)를 가진 문단을 모두 제목으로
            for sh, para, raw, fp in paras:
                if fp is not None and abs(fp - max_pt) < 0.01:
                    title_ids.add(id(para))
        elif title_shape is not None:
            # 크기 정보가 전무 → title placeholder shape의 문단을 제목으로
            # (paras에 담긴 것과 같은 shape인지로 판단; para 재읽기는 id가 달라짐)
            for sh, para, raw, fp in paras:
                if sh == title_shape:
                    title_ids.add(id(para))

        # 3) 라인 생성 (제목은 ##, 나머지는 볼드/불릿 유지)
        title_lines, body_lines = [], []
        for sh, para, raw, fp in paras:
            if id(para) in title_ids:
                # 제목: 볼드 마커 없이 순수 텍스트에 "## "
                title_lines.append(f"## {raw.strip()}")
            else:
                body = _runs_to_markdown(para.runs).strip()
                if getattr(para, "level", 0) and para.level > 0:
                    body_lines.append(f"- {body}")
                else:
                    body_lines.append(body)
        slides.append("\n".join(title_lines + body_lines))
    return slides, tables


def extract_pptx(path: str):
    """슬라이드를 순서대로 이어붙인 flat 버전(본문 텍스트, 표 리스트)."""
    slides, tables = extract_pptx_slides(path)
    return "\n".join(s for s in slides if s.strip()), tables


def _pptx_table_to_html(table) -> str:
    rows_html = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append(f"<td>{html_lib.escape(cell.text.strip())}</td>")
        rows_html.append(f"<tr>{''.join(cells)}</tr>")
    return f"<table>{''.join(rows_html)}</table>"


# ══════════════════════════════════════════════════════════════════════
# .xlsx 추출
# ══════════════════════════════════════════════════════════════════════
def extract_xlsx(path: str, include_hidden: bool = True, skip_hidden_sheets: bool = True):
    """
    엑셀은 시트 자체가 표이므로, 각 시트를 <table> 하나로 변환한다.
    본문 text에는 각 시트명을 "## 시트명" 형태로 (표 위 헤더처럼) 담아 반환한다.
    파서도 시트명을 "## 헤더" + 표 형태로 내므로 표기 체계를 맞춘 것이다.

    처리 규칙
    ---------
      - 빈 셀(None)은 빈 문자열로, 나머지는 문자열로 변환하고 바깥쪽 빈 행/열은 잘라낸다.
      - 병합 셀은 좌상단 값을 병합 범위 전체에 채운다(파서가 병합을 어떻게 풀든
        내용 비교가 되도록).
      - 수식 셀은 저장된 계산값(data_only)을 읽는다. 엑셀이 값을 캐시하지 않았으면
        None이 될 수 있다.
      - skip_hidden_sheets=True면 숨긴 시트는 건너뛴다.
      - include_hidden=False면 숨긴 행/열도 제외한다(필터/숨기기로 가려진 내용을
        파서가 안 읽는다고 볼 때 공정 비교용). 기본값 True는 전부 포함.
    (차트/이미지/피벗 등 셀 밖 객체는 추출하지 않는다.)
    """
    from openpyxl import load_workbook
    from openpyxl.utils import range_boundaries, column_index_from_string

    # 병합/숨김 정보를 읽으려면 read_only=False 여야 한다(대용량이면 메모리 주의).
    wb = load_workbook(path, data_only=True)
    tables = []
    text_parts = []   # 시트명(## ...)을 등장 순서대로 모아 본문 텍스트로 반환

    for ws in wb.worksheets:
        if skip_hidden_sheets and ws.sheet_state != "visible":
            continue

        n_rows, n_cols = ws.max_row or 0, ws.max_column or 0
        if n_rows == 0 or n_cols == 0:
            continue

        grid = [["" for _ in range(n_cols)] for _ in range(n_rows)]
        for row in ws.iter_rows():
            for cell in row:
                v = cell.value
                grid[cell.row - 1][cell.column - 1] = "" if v is None else str(v).strip()

        # 병합 셀: 좌상단 값을 범위 전체에 채움
        for merged in ws.merged_cells.ranges:
            min_c, min_r, max_c, max_r = range_boundaries(str(merged))
            top_left = grid[min_r - 1][min_c - 1]
            for r in range(min_r, max_r + 1):
                for c in range(min_c, max_c + 1):
                    grid[r - 1][c - 1] = top_left

        # 숨긴 행/열 제외 (옵션)
        if not include_hidden:
            hidden_rows = {i for i, d in ws.row_dimensions.items() if d.hidden}
            hidden_cols = set()
            for letter, d in ws.column_dimensions.items():
                if d.hidden:
                    try:
                        hidden_cols.add(column_index_from_string(letter))
                    except ValueError:
                        pass
            grid = [
                [v for ci, v in enumerate(row, start=1) if ci not in hidden_cols]
                for ri, row in enumerate(grid, start=1) if ri not in hidden_rows
            ]

        matrix = _trim_empty_edges(grid)
        if matrix:
            # 시트명을 "## 시트명"으로 본문에 넣는다(표 위에 한 줄 띄운다).
            # 파서 출력 예: "## 헤더구조,음영추출\n\n| null | null | ... |" 형태.
            text_parts.append(f"## {ws.title}")
            tables.append(_matrix_to_html(matrix))

    wb.close()
    # 시트명들 사이는 빈 줄로 구분(파서가 시트별로 ## 헤더 + 표를 내는 것과 맞춤)
    return "\n\n".join(text_parts), tables


def _trim_empty_edges(matrix):
    """행렬 바깥쪽의 완전히 빈 행/열을 제거한다."""
    # 빈 행 제거
    matrix = [row for row in matrix if any(c != "" for c in row)]
    if not matrix:
        return []
    n_cols = max(len(r) for r in matrix)
    matrix = [r + [""] * (n_cols - len(r)) for r in matrix]  # 길이 맞춤
    # 완전히 빈 열 제거
    keep_cols = [c for c in range(n_cols) if any(row[c] != "" for row in matrix)]
    return [[row[c] for c in keep_cols] for row in matrix]


def _matrix_to_html(matrix) -> str:
    rows_html = []
    for row in matrix:
        cells = "".join(f"<td>{html_lib.escape(v)}</td>" for v in row)
        rows_html.append(f"<tr>{cells}</tr>")
    return f"<table>{''.join(rows_html)}</table>"


# ══════════════════════════════════════════════════════════════════════
# 확장자 기준 자동 분기
# ══════════════════════════════════════════════════════════════════════
_EXTRACTORS = {".docx": extract_docx, ".pptx": extract_pptx, ".xlsx": extract_xlsx}


def extract_gt(path: str):
    """
    확장자를 보고 알맞은 추출기를 골라 (text, tables)를 반환한다.
    지원: .docx / .pptx / .xlsx
    """
    ext = Path(path).suffix.lower()
    if ext not in _EXTRACTORS:
        raise ValueError(f"지원하지 않는 GT 형식입니다: {ext} (지원: .docx/.pptx/.xlsx)")
    return _EXTRACTORS[ext](path)