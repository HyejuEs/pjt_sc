# -*- coding: utf-8 -*-
"""
build_test_suite.py
====================
가이드라인(공통/텍스트/표/이미지/Word/Excel/Powerpoint)에 나온 "문제사항 → 조치사항"을
실제로 재현한 테스트 문서 세트를 만든다.

핵심 설계: "듀얼 GT"
--------------------
가이드라인 적용은 파서가 아니라 "원문 문서 자체"를 고치는 작업이다. 그래서 데이터셋도
문서id마다 아래 5개 파일을 만든다.

    {id}_before_gt.(docx|pptx|xlsx)   가이드라인 적용 "전" 원문 (문제 재현)
    {id}_after_gt.(docx|pptx|xlsx)    가이드라인 적용 "후" 원문 (조치사항 반영, 확장자 동일)
    {id}_before.txt                   before_gt를 "실제 파서"가 읽었다고 가정한 결과(시뮬레이션)
    {id}_after.txt                    after_gt를 그 파서가 읽었다고 가정한 결과(시뮬레이션)

before.txt는 before_gt와, after.txt는 after_gt와 비교된다(run_eval.py의 듀얼 GT 지원).
이렇게 해야 "원문 개선 + 파싱 난이도 감소"의 최종 효과를 재는, 사용자가 원한 비교가 된다.

파서 결과는 실제 파서가 없으므로 "이 문제 유형이면 파서가 전형적으로 이렇게 망가진다"는
패턴을 손으로 흉내낸 시뮬레이션이다(각 함수 docstring에 어떤 손상을 재현했는지 명시).
"""

from pathlib import Path
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt, Emu
from openpyxl import Workbook
from openpyxl.styles import PatternFill

OUT = Path("./testset")
OUT.mkdir(exist_ok=True)


# ══════════════════════════════════════════════════════════════════════
# docx 유틸
# ══════════════════════════════════════════════════════════════════════
def _add_textbox(paragraph, text):
    """VML 텍스트박스를 문단에 삽입(python-docx엔 고수준 API가 없어 raw XML 사용)."""
    run = paragraph.add_run()
    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
          "v": "urn:schemas-microsoft-com:vml"}
    xml = (f'<w:pict xmlns:w="{ns["w"]}" xmlns:v="{ns["v"]}">'
           f'<v:shape style="width:220pt;height:70pt">'
           f'<v:textbox><w:txbxContent><w:p><w:r><w:t>{text}</w:t></w:r></w:p></w:txbxContent></v:textbox>'
           f'</v:shape></w:pict>')
    run._r.append(etree.fromstring(xml.encode("utf-8")))


def _strike_run(run):
    run.font.strike = True


def _set_cell_shading(cell, color="FFFF00"):
    """셀 배경 하이라이트(볼드가 아닌 음영 방식 강조 - '셀 강조' 문제 재현용)."""
    tcPr = cell._tc.get_or_add_tcPr()
    shd = etree.SubElement(tcPr, qn("w:shd"))
    shd.set(qn("w:val"), "clear"); shd.set(qn("w:fill"), color)


# ══════════════════════════════════════════════════════════════════════
# [DOC 1] word_text_table : 텍스트/표 카테고리 - docx
#   담는 문제: 제목 없음, 소제목 중복, 특수문자 미정의, 약어 미정의, 한자 미병기,
#             OX/YN 기호, 텍스트 강조(하이라이트 vs 볼드), 표 복잡헤더(2단 초과),
#             표 동일헤더 반복, 표 하단헤더, 표 불필요 셀병합
# ══════════════════════════════════════════════════════════════════════
def make_word_text_table():
    doc_id = "word_text_table"

    # ── before_gt: 문제 재현 ─────────────────────────────────────────
    d = Document()
    # 제목 없음: 헤딩 스타일 없이 바로 본문 시작(섹션 제목 누락)
    d.add_paragraph("신상품 광고물 사전심의는 마케팅팀 담당자가 접수 후 3영업일 이내 처리한다.")

    # 소제목 중복: 서로 다른 내용인데 같은 소제목("세부절차")을 반복
    d.add_heading("세부절차", level=2)
    d.add_paragraph("① 접수 ② 검토 ③ 승인")
    d.add_heading("세부절차", level=2)  # 중복 - 실제로는 "특별점검 세부절차"여야 함
    d.add_paragraph("① 현장점검 ② 결과보고 ③ 시정조치")

    # 특수문자 미정의 + 약어 미정의 + 한자 미병기 + OX 비표준
    p = d.add_paragraph("PF★ 대출은 승인 대상에서 제외되며(★=고위험군), 준법감시인(CCO)이 최종 확인한다.")
    d.add_paragraph("동 규정은 甲(발행사)과 乙(대행사) 간 계약에 적용된다.")  # 한자 병기 없음
    d.add_paragraph("적합성 원칙 준수 여부: O   설명의무 이행 여부: X")  # OX 표기, 범례 없음

    # 텍스트 강조: 볼드 대신 하이라이트(음영)만 사용 - run에 highlight 색만 주고 bold=False
    p2 = d.add_paragraph()
    r = p2.add_run("반드시 준수해야 하는 핵심 조항")
    r.font.highlight_color = 7  # WD_COLOR_INDEX.YELLOW - bold는 아님
    r.bold = False

    # 표: 2단 초과 복잡 헤더 + 동일 헤더 반복 + 불필요 셀병합 + 하단헤더
    d.add_heading("월별 대출 실적 비교", level=1)
    t = d.add_table(rows=5, cols=4)
    t.style = "Table Grid"
    # 3단 헤더(구분/2025년 상반기/2025년 하반기 - 상반기 밑에 또 (1월,2월) 등 => 여기선 문자열로 흉내)
    t.rows[0].cells[0].text = "구분"
    t.rows[0].cells[1].text = "2025 상반기"
    t.rows[0].cells[2].text = "2025 상반기"  # 동일 헤더 반복(병합 대신 텍스트 중복)
    t.rows[0].cells[3].text = "비고"
    t.rows[1].cells[0].text = "구분"          # 헤더가 또 나옴(동일 헤더 반복)
    t.rows[1].cells[1].text = "1월"
    t.rows[1].cells[2].text = "2월"
    t.rows[1].cells[3].text = ""
    t.rows[2].cells[0].text = "저축은행"
    t.rows[2].cells[1].text = "429"
    t.rows[2].cells[2].text = "322"
    t.rows[2].cells[3].text = "전월대비 감소"
    t.rows[3].cells[0].text = "인터넷은행"
    t.rows[3].cells[1].text = "3329"
    t.rows[3].cells[2].text = "2605"
    t.rows[3].cells[3].text = "전월대비 감소"
    # 하단헤더: 표 맨 아래 행에 헤더성 캡션을 넣음(표 상단엔 제목 없음)
    t.rows[4].cells[0].merge(t.rows[4].cells[3]).text = "[표] 저축은행/인터넷은행 월별 대출 실적(단위: 억원)"
    d.save(OUT / f"{doc_id}_before_gt.docx")

    # ── after_gt: 조치사항 반영 ──────────────────────────────────────
    d2 = Document()
    d2.add_heading("신상품 광고물 사전심의 절차", level=1)
    d2.add_paragraph("신상품 광고물 사전심의는 마케팅팀 담당자가 접수 후 3영업일 이내 처리한다.")

    d2.add_heading("일반 심의 세부절차", level=2)
    d2.add_paragraph("① 접수 ② 검토 ③ 승인")
    d2.add_heading("특별점검 세부절차", level=2)  # 고유 제목으로 구분
    d2.add_paragraph("① 현장점검 ② 결과보고 ③ 시정조치")

    d2.add_paragraph(
        "PF(프로젝트 파이낸싱)★ 대출은 승인 대상에서 제외되며(★: 고위험군을 의미), "
        "준법감시인(CCO, Chief Compliance Officer)이 최종 확인한다."
    )
    d2.add_paragraph("동 규정은 갑(甲, 발행사)과 을(乙, 대행사) 간 계약에 적용된다.")
    d2.add_paragraph("적합성 원칙 준수 여부: O(준수)   설명의무 이행 여부: X(미이행)")

    p2b = d2.add_paragraph()
    r2 = p2b.add_run("반드시 준수해야 하는 핵심 조항")
    r2.bold = True  # 강조는 볼드로 표준화

    d2.add_heading("월별 대출 실적 비교", level=1)  # 표 위에 표준 제목
    t2 = d2.add_table(rows=3, cols=4)
    t2.style = "Table Grid"
    t2.rows[0].cells[0].text = "구분"
    t2.rows[0].cells[1].text = "2025년 상반기 1월(억원)"
    t2.rows[0].cells[2].text = "2025년 상반기 2월(억원)"
    t2.rows[0].cells[3].text = "비고"
    t2.rows[1].cells[0].text = "저축은행"
    t2.rows[1].cells[1].text = "429"
    t2.rows[1].cells[2].text = "322"
    t2.rows[1].cells[3].text = "전월대비 감소"
    t2.rows[2].cells[0].text = "인터넷은행"
    t2.rows[2].cells[1].text = "3329"
    t2.rows[2].cells[2].text = "2605"
    t2.rows[2].cells[3].text = "전월대비 감소"
    d2.save(OUT / f"{doc_id}_after_gt.docx")

    # ── 파서 결과 시뮬레이션 ─────────────────────────────────────────
    # before: 제목 없어 문맥 인식 실패로 앞문장 누락 처리되진 않지만(텍스트 자체는 인식),
    #         소제목 중복은 파서가 두 섹션을 구분 못해 잘못 병합, 특수문자/한자/OX는
    #         파서가 다르게 표기(★→*, 甲/乙 못읽음, O/X 그대로), highlight는 손실(**없음),
    #         표는 헤더 2행을 그대로 살려 pandas류 파서처럼 "Unnamed" 흔적 + 하단 캡션 인식 실패.
    before_txt = (
        "신상품 광고물 사전심의는 마케팅팀 담당자가 접수 후 3영업일 이내 처리한다.\n"
        "## 세부절차\n① 접수 ② 검토 ③ 승인\n"
        "## 세부절차\n① 현장점검 ② 결과보고 ③ 시정조치\n"  # 제목 구분 안 됨(둘 다 "세부절차")
        "PF* 대출은 승인 대상에서 제외되며, 준법감시인이 최종 확인한다.\n"  # ★→*, CCO 누락
        "동 규정은 발행사와 대행사 간 계약에 적용된다.\n"  # 甲/乙 누락(한자 인식 실패)
        "적합성 원칙 준수 여부: O   설명의무 이행 여부: X\n"
        "반드시 준수해야 하는 핵심 조항\n"  # 볼드 기호 없음(하이라이트 손실)
        # 파서가 표 셀 내용을 본문으로도 한 번 더 쏟아낸 중복(표중복 제거 기능이 걸러냄)
        "저축은행 429 322\n"
        "인터넷은행 3329 2605\n"
        "<table><tr><td>구분</td><td>2025 상반기</td><td>2025 상반기</td><td>비고</td></tr>"
        "<tr><td>구분</td><td>1월</td><td>2월</td><td></td></tr>"
        "<tr><td>저축은행</td><td>429</td><td>322</td><td>전월대비 감소</td></tr>"
        "<tr><td>인터넷은행</td><td>3329</td><td>2605</td><td>전월대비 감소</td></tr>"
        "<tr><td>[표] 저축은행/인터넷은행 월별 대출 실적(단위: 억원)</td><td></td><td></td><td></td></tr></table>\n"
    )
    (OUT / f"{doc_id}_before.txt").write_text(before_txt, encoding="utf-8")

    # after: 표준화된 원문이라 파서가 대부분 정확히 인식(약간의 강조 유실 정도만 잔존)
    after_txt = (
        "# 신상품 광고물 사전심의 절차\n"
        # 파서가 GT 한 문장을 두 줄로 쪼갬 → 줄바꿈 유연화가 GT 한 줄로 합쳐 WER에 반영 안 함
        "신상품 광고물 사전심의는 마케팅팀 담당자가\n접수 후 3영업일 이내 처리한다.\n"
        "## 일반 심의 세부절차\n① 접수 ② 검토 ③ 승인\n"
        "## 특별점검 세부절차\n① 현장점검 ② 결과보고 ③ 시정조치\n"
        "PF(프로젝트 파이낸싱)★ 대출은 승인 대상에서 제외되며(★: 고위험군을 의미), "
        "준법감시인(CCO, Chief Compliance Officer)이 최종 확인한다.\n"
        "동 규정은 갑(甲, 발행사)과 을(乙, 대행사) 간 계약에 적용된다.\n"
        "적합성 원칙 준수 여부: O(준수)   설명의무 이행 여부: X(미이행)\n"
        "**반드시 준수해야 하는 핵심 조항**\n"
        "# 월별 대출 실적 비교\n"
        "<table><tr><td>구분</td><td>2025년 상반기 1월(억원)</td><td>2025년 상반기 2월(억원)</td><td>비고</td></tr>"
        "<tr><td>저축은행</td><td>429</td><td>322</td><td>전월대비 감소</td></tr>"
        "<tr><td>인터넷은행</td><td>3329</td><td>2605</td><td>전월대비 감소</td></tr></table>\n"
    )
    (OUT / f"{doc_id}_after.txt").write_text(after_txt, encoding="utf-8")
    print(f"[생성] {doc_id}")


# ══════════════════════════════════════════════════════════════════════
# [DOC 2] word_box_memo_strike : 텍스트박스 / 메모 / 취소선 - docx
# ══════════════════════════════════════════════════════════════════════
def make_word_box_memo_strike():
    doc_id = "word_box_memo_strike"

    d = Document()
    d.add_heading("모집인 자격 정지 기준 변경", level=1)
    p = d.add_paragraph("자격 정지 기준은 아래와 같다: ")
    _add_textbox(p, "예외: 신규 입사 6개월 이내는 정지 대신 경고 처리")  # 핵심 예외가 텍스트박스에만 있음

    p2 = d.add_paragraph("기존 기준 3회 위반 시 정지였으나, ")
    r_old = p2.add_run("이번 개정으로 2회 위반 시 정지")
    _strike_run(r_old)  # 취소선(1): 개정 전/후 비교 정보를 취소선으로만 표현
    p2.add_run("로 변경한다.")

    p3 = d.add_paragraph("담당 부서는 ")
    r_del = p3.add_run("영업지원팀이었으나 준법지원팀으로 이관됨")
    _strike_run(r_del)  # 취소선(2): 보존 불필요한 옛 내용 삭제 대상
    d.add_comment([p3.runs[0]], text="영업지원팀 언급은 실수 - 원래 준법지원팀이 맞음. 본문에도 반영 필요", author="검토자A")
    d.save(OUT / f"{doc_id}_before_gt.docx")

    d2 = Document()
    d2.add_heading("모집인 자격 정지 기준 변경", level=1)
    d2.add_paragraph(
        "자격 정지 기준은 아래와 같다: 신규 입사 6개월 이내는 정지 대신 경고 처리한다."
    )  # 텍스트박스 내용을 본문에 통합
    d2.add_paragraph("기존 기준은 3회 위반 시 정지였으나, 이번 개정으로 2회 위반 시 정지로 변경한다.")
    # 변경 전 내용은 비교 목적상 보존 -> 표로 정리
    t = d2.add_table(rows=2, cols=2)
    t.style = "Table Grid"
    t.rows[0].cells[0].text = "개정 전"; t.rows[0].cells[1].text = "개정 후"
    t.rows[1].cells[0].text = "3회 위반 시 정지"; t.rows[1].cells[1].text = "2회 위반 시 정지"
    d2.add_paragraph("담당 부서는 준법지원팀이다.")  # 취소선/메모 내용 정리해 본문에 확정 반영, 메모 제거
    d2.save(OUT / f"{doc_id}_after_gt.docx")

    before_txt = (
        "모집인 자격 정지 기준 변경\n"
        "자격 정지 기준은 아래와 같다: \n"  # 텍스트박스 내용 통째로 누락(도형/텍스트박스는 파서가 못 읽음)
        "기존 기준 3회 위반 시 정지였으나, 로 변경한다.\n"  # 취소선 내용도 통째로 누락
        "담당 부서는 이다.\n"  # 취소선 내용 누락 + 메모는 원래 안 읽힘
    )
    (OUT / f"{doc_id}_before.txt").write_text(before_txt, encoding="utf-8")

    after_txt = (
        "모집인 자격 정지 기준 변경\n"
        "자격 정지 기준은 아래와 같다: 신규 입사 6개월 이내는 정지 대신 경고 처리한다.\n"
        "기존 기준은 3회 위반 시 정지였으나, 이번 개정으로 2회 위반 시 정지로 변경한다.\n"
        "<table><tr><td>개정 전</td><td>개정 후</td></tr>"
        "<tr><td>3회 위반 시 정지</td><td>2회 위반 시 정지</td></tr></table>\n"
        "담당 부서는 준법지원팀이다.\n"
    )
    (OUT / f"{doc_id}_after.txt").write_text(after_txt, encoding="utf-8")
    print(f"[생성] {doc_id}  (텍스트박스/메모 GT 포함 여부는 gt_extract 확장으로 검증 가능)")


# ══════════════════════════════════════════════════════════════════════
# [DOC 3] ppt_slides : 슬라이드 경계 밖 콘텐츠 / 도형표 / 다이어그램 / 제목중복 - pptx
# ══════════════════════════════════════════════════════════════════════
def make_ppt_slides():
    doc_id = "ppt_slides"

    def _add_slide(prs, title):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        return slide

    # ── before ──
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)

    s1 = _add_slide(prs, "JURIS 준법심의 절차")
    tb = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2)).text_frame
    tb.text = "심의는 접수 후 진행한다"
    # 슬라이드 경계 밖에 핵심 정보 배치(문제)
    tb2 = s1.shapes.add_textbox(Inches(9.5), Inches(6), Inches(4), Inches(1)).text_frame
    tb2.text = "예외: 긴급 광고물은 사전 승인 없이 사후 심의 가능"

    s2 = _add_slide(prs, "특별점검 절차")  # 이후 제목 중복될 슬라이드
    tf2 = s2.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3)).text_frame
    tf2.text = "정기점검과 특별점검으로 나뉜다"
    p = tf2.add_paragraph(); p.text = "현장점검은 실시한다"  # 순서/빈도 정보 누락("분기 1회" 빠짐)

    # 도형 조합 표(실제 표 대신 텍스트박스 격자로 흉내) - 파서는 개별 텍스트만 인식, 구조 인식 실패
    s3 = _add_slide(prs, "특별점검 절차")  # 제목 중복(서로 다른 내용인데 동일 제목)
    labels = [("구분", 1), ("1월", 3), ("2월", 5), ("저축은행", 1), ("429", 3), ("322", 5)]
    y = 2
    for i, (txt, x) in enumerate(labels):
        yy = 2 if i < 3 else 3
        box = s3.shapes.add_textbox(Inches(x), Inches(yy), Inches(1.5), Inches(0.6))
        box.text_frame.text = txt
    prs.save(OUT / f"{doc_id}_before_gt.pptx")

    # ── after ──
    prs2 = Presentation()
    prs2.slide_width = Inches(10); prs2.slide_height = Inches(7.5)

    s1b = _add_slide(prs2, "JURIS 준법심의 절차")
    tfb = s1b.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2)).text_frame
    tfb.text = "심의는 접수 후 진행한다"
    p = tfb.add_paragraph()
    p.text = "예외: 긴급 광고물은 사전 승인 없이 사후 심의 가능"  # 경계 안으로 재배치
    p.level = 1

    s2b = _add_slide(prs2, "특별점검 절차")
    tf2b = s2b.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3)).text_frame
    tf2b.text = "정기점검과 특별점검으로 나뉜다"
    p2b = tf2b.add_paragraph(); p2b.text = "현장점검은 분기 1회 실시한다"  # 순서/빈도 정보 보강

    s3b = _add_slide(prs2, "특별점검 결과표")  # 고유한 제목
    table_shape = s3b.shapes.add_table(2, 3, Inches(1), Inches(2), Inches(6), Inches(1.2))
    tbl = table_shape.table
    tbl.cell(0, 0).text = "구분"; tbl.cell(0, 1).text = "1월(억원)"; tbl.cell(0, 2).text = "2월(억원)"
    tbl.cell(1, 0).text = "저축은행"; tbl.cell(1, 1).text = "429"; tbl.cell(1, 2).text = "322"
    prs2.save(OUT / f"{doc_id}_after_gt.pptx")

    before_txt = (
        "## JURIS 준법심의 절차\n심의는 접수 후 진행한다\n"
        # 슬라이드 경계 밖 텍스트박스는 파서가 아예 못 읽음(누락)
        "## 특별점검 절차\n정기점검과 특별점검으로 나뉜다\n- 현장점검은 실시한다\n"
        "## 특별점검 절차\n"  # 이전 슬라이드와 제목 동일 -> 파서가 이후 섹션 구분/정렬에 혼란
        "구분\n1월\n2월\n저축은행\n429\n322\n"  # 도형 조합 표 -> 표 구조 인식 실패, 순서 뒤섞임
    )
    (OUT / f"{doc_id}_before.txt").write_text(before_txt, encoding="utf-8")

    after_txt = (
        "## JURIS 준법심의 절차\n심의는 접수 후 진행한다\n- 예외: 긴급 광고물은 사전 승인 없이 사후 심의 가능\n"
        "## 특별점검 절차\n정기점검과 특별점검으로 나뉜다\n- 현장점검은 분기 1회 실시한다\n"
        "## 특별점검 결과표\n"
        "<table><tr><td>구분</td><td>1월(억원)</td><td>2월(억원)</td></tr>"
        "<tr><td>저축은행</td><td>429</td><td>322</td></tr></table>\n"
    )
    (OUT / f"{doc_id}_after.txt").write_text(after_txt, encoding="utf-8")
    print(f"[생성] {doc_id}")


# ══════════════════════════════════════════════════════════════════════
# [DOC 4] xlsx_single : 숨김 행/열, 드롭다운, 병합, 하단헤더 - xlsx
# ══════════════════════════════════════════════════════════════════════
def make_xlsx_single():
    doc_id = "xlsx_single"

    wb = Workbook(); ws = wb.active; ws.title = "실적"
    ws.append(["구분", "1월", "2월", "내부메모"])
    ws.append(["저축은행", 429, 322, "재검토 필요"])
    ws.append(["인터넷은행", 3329, 2605, "정상"])
    ws.column_dimensions["D"].hidden = True  # 숨김 열에 핵심 아닌 메모(그나마 안전한 예)
    ws.append(["카드사", 900, 810, ""])
    ws.row_dimensions[4].hidden = True  # 숨김 행에 실제 데이터가 숨어 있음(문제)
    ws.auto_filter.ref = "A1:D3"  # 드롭다운(필터) 사용
    wb.save(OUT / f"{doc_id}_before_gt.xlsx")

    wb2 = Workbook(); ws2 = wb2.active; ws2.title = "실적"
    ws2.append(["구분", "1월(억원)", "2월(억원)"])
    ws2.append(["저축은행", 429, 322])
    ws2.append(["인터넷은행", 3329, 2605])
    ws2.append(["카드사", 900, 810])  # 숨김 해제, 내부메모 열은 삭제(불필요 참조 제거)
    wb2.save(OUT / f"{doc_id}_after_gt.xlsx")

    # before: 숨김 행/열은 파서가 기본적으로 건너뜀 -> 카드사 행, 내부메모 열 통째로 누락
    before_txt = (
        "## 실적\n"
        "| 구분 | 1월 | 2월 |\n"
        "| 저축은행 | 429 | 322 |\n"
        "| 인터넷은행 | 3329 | 2605 |\n"
    )
    (OUT / f"{doc_id}_before.txt").write_text(before_txt, encoding="utf-8")

    after_txt = (
        "## 실적\n"
        "| 구분 | 1월(억원) | 2월(억원) |\n"
        "| 저축은행 | 429 | 322 |\n"
        "| 인터넷은행 | 3329 | 2605 |\n"
        "| 카드사 | 900 | 810 |\n"
    )
    (OUT / f"{doc_id}_after.txt").write_text(after_txt, encoding="utf-8")
    print(f"[생성] {doc_id}")


# ══════════════════════════════════════════════════════════════════════
# [DOC 5] xlsx_multisheet_linked : 여러 시트가 유기적으로 연결된 케이스 (요청사항)
#   시트1 "지점별_실적_요약", 시트2 "서울지점_상세", 시트3 "부산지점_상세"가
#   같은 분기 실적이라는 주제로 묶여 있는데, before는 그 연결관계가 텍스트로
#   전혀 명시되지 않아(시트명 외엔 단서 없음) 파서가 시트들을 별개로 처리한다.
#   after는 가이드라인("연관정보는 헤더/캡션으로 명시")에 따라 각 시트 상단에
#   연결 캡션을 넣는다.
# ══════════════════════════════════════════════════════════════════════
def make_xlsx_multisheet_linked():
    doc_id = "xlsx_multisheet_linked"

    wb = Workbook()
    ws1 = wb.active; ws1.title = "지점별_실적_요약"
    ws1.append(["지점", "3분기 매출", "3분기 비용"])
    ws1.append(["서울", 5200, 3100])
    ws1.append(["부산", 3100, 2000])

    ws2 = wb.create_sheet("서울지점_상세")
    ws2.append(["월", "매출", "비용"])  # 이 표가 요약시트 "서울" 행의 근거라는 언급이 전혀 없음
    ws2.append(["7월", 1700, 1000])
    ws2.append(["8월", 1800, 1050])
    ws2.append(["9월", 1700, 1050])

    ws3 = wb.create_sheet("부산지점_상세")
    ws3.append(["월", "매출", "비용"])
    ws3.append(["7월", 1000, 650])
    ws3.append(["8월", 1050, 680])
    ws3.append(["9월", 1050, 670])
    wb.save(OUT / f"{doc_id}_before_gt.xlsx")

    wb2 = Workbook()
    wsa = wb2.active; wsa.title = "지점별_실적_요약"
    wsa.append(["[3분기 지점 실적 요약 - 상세는 서울지점_상세/부산지점_상세 시트 참조]"])
    wsa.append(["지점", "3분기 매출(백만원)", "3분기 비용(백만원)"])
    wsa.append(["서울", 5200, 3100])
    wsa.append(["부산", 3100, 2000])

    wsb = wb2.create_sheet("서울지점_상세")
    wsb.append(["[지점별_실적_요약 시트 '서울' 행의 월별 상세 내역]"])
    wsb.append(["월", "매출(백만원)", "비용(백만원)"])
    wsb.append(["7월", 1700, 1000])
    wsb.append(["8월", 1800, 1050])
    wsb.append(["9월", 1700, 1050])

    wsc = wb2.create_sheet("부산지점_상세")
    wsc.append(["[지점별_실적_요약 시트 '부산' 행의 월별 상세 내역]"])
    wsc.append(["월", "매출(백만원)", "비용(백만원)"])
    wsc.append(["7월", 1000, 650])
    wsc.append(["8월", 1050, 680])
    wsc.append(["9월", 1050, 670])
    wb2.save(OUT / f"{doc_id}_after_gt.xlsx")

    # 파서는 시트명을 "## 시트명"으로 내보낸다(GT 추출과 동일 체계).
    # before: 시트명만 있고 시트 간 "연결 정보" 텍스트가 없음(문제 상황).
    before_txt = (
        "## 지점별_실적_요약\n"
        "| 지점 | 3분기 매출 | 3분기 비용 |\n| 서울 | 5200 | 3100 |\n| 부산 | 3100 | 2000 |\n\n"
        "## 서울지점_상세\n"
        "| 월 | 매출 | 비용 |\n| 7월 | 1700 | 1000 |\n| 8월 | 1800 | 1050 |\n| 9월 | 1700 | 1050 |\n\n"
        "## 부산지점_상세\n"
        "| 월 | 매출 | 비용 |\n| 7월 | 1000 | 650 |\n| 8월 | 1050 | 680 |\n| 9월 | 1050 | 670 |\n"
    )
    (OUT / f"{doc_id}_before.txt").write_text(before_txt, encoding="utf-8")

    # after: 각 시트 상단(A1 셀)에 연결 캡션이 들어간다. xlsx는 모든 게 셀이므로
    #   캡션도 "표의 첫 행"이다 → before_gt엔 없고 after_gt엔 있는 이 캡션 행 차이가
    #   TEDS(표 구조/내용)로 잡힌다. 시트명(##)은 before/after 동일하라 WER은 0.
    after_txt = (
        "## 지점별_실적_요약\n"
        "| 3분기 지점 실적 요약 - 상세는 서울지점_상세/부산지점_상세 시트 참조 | | |\n"
        "| 지점 | 3분기 매출(백만원) | 3분기 비용(백만원) |\n| 서울 | 5200 | 3100 |\n| 부산 | 3100 | 2000 |\n\n"
        "## 서울지점_상세\n"
        "| 지점별_실적_요약 시트 서울 행의 월별 상세 내역 | | |\n"
        "| 월 | 매출(백만원) | 비용(백만원) |\n| 7월 | 1700 | 1000 |\n| 8월 | 1800 | 1050 |\n| 9월 | 1700 | 1050 |\n\n"
        "## 부산지점_상세\n"
        "| 지점별_실적_요약 시트 부산 행의 월별 상세 내역 | | |\n"
        "| 월 | 매출(백만원) | 비용(백만원) |\n| 7월 | 1000 | 650 |\n| 8월 | 1050 | 680 |\n| 9월 | 1050 | 670 |\n"
    )
    (OUT / f"{doc_id}_after.txt").write_text(after_txt, encoding="utf-8")
    print(f"[생성] {doc_id}  (다중 시트 연결 케이스)")


# ══════════════════════════════════════════════════════════════════════
# [DOC 6] ppt_table_dup : 표 밖 중복 텍스트 제거(WER만) 검증용 - pptx
#   before: 파서가 표를 <table>로도 내고, 같은 셀 내용을 표 밖 본문으로도 중복 추출.
#           → WER 계산 시 그 중복 본문 라인이 제거되어야 하고(삽입 오류로 안 잡힘),
#             리포트에는 "표 내용 중복 추출"로 따로 표시되어야 한다.
#   after : 중복 없이 표만 깔끔히. (TEDS는 두 경우 다 표를 그대로 보므로 영향 없음)
# ══════════════════════════════════════════════════════════════════════
def make_ppt_table_dup():
    doc_id = "ppt_table_dup"

    def _add_title_slide(prs, title):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        return slide

    # GT(before/after 공통 원본 개념이지만, 듀얼 GT 규약을 위해 둘 다 생성)
    def _build_gt(path):
        prs = Presentation()
        prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
        s = _add_title_slide(prs, "분기 실적 요약")
        tf = s.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1)).text_frame
        tf.text = "아래 표는 은행별 분기 실적이다"
        table_shape = s.shapes.add_table(3, 3, Inches(1), Inches(3), Inches(6), Inches(2))
        t = table_shape.table
        t.cell(0, 0).text = "구분"; t.cell(0, 1).text = "1분기"; t.cell(0, 2).text = "2분기"
        t.cell(1, 0).text = "저축은행"; t.cell(1, 1).text = "429"; t.cell(1, 2).text = "322"
        t.cell(2, 0).text = "인터넷은행"; t.cell(2, 1).text = "3329"; t.cell(2, 2).text = "2605"
        prs.save(path)

    _build_gt(OUT / f"{doc_id}_before_gt.pptx")
    _build_gt(OUT / f"{doc_id}_after_gt.pptx")

    # before/after 상관없이: txt 안에서 표(HTML)의 값이 표 밖 본문에도 통째로 중복 추출된
    # 상황. WER은 이 중복을 삭제/수정 없이 그대로 포함해 계산하고, report.html에서만
    # 그 중복 라인을 파란색으로 '표시'한다. (여기선 before/after 둘 다 같은 중복 상황을 둠)
    dup_body = (
        "구분 1분기 2분기\n"       # ← 표 내용이 본문에도 통째로 중복 추출됨(표시 대상)
        "저축은행 429 322\n"
        "인터넷은행 3329 2605\n"
    )
    table_html = (
        "<table><tr><td>구분</td><td>1분기</td><td>2분기</td></tr>"
        "<tr><td>저축은행</td><td>429</td><td>322</td></tr>"
        "<tr><td>인터넷은행</td><td>3329</td><td>2605</td></tr></table>\n"
    )
    before_txt = "## 분기 실적 요약\n아래 표는 은행별 분기 실적이다\n" + dup_body + table_html
    after_txt = "## 분기 실적 요약\n아래 표는 은행별 분기 실적이다\n" + dup_body + table_html
    (OUT / f"{doc_id}_before.txt").write_text(before_txt, encoding="utf-8")
    (OUT / f"{doc_id}_after.txt").write_text(after_txt, encoding="utf-8")
    print(f"[생성] {doc_id}  (표 밖 중복 텍스트: WER 포함 + 리포트 표시)")


if __name__ == "__main__":
    make_word_text_table()
    make_word_box_memo_strike()
    make_ppt_slides()
    make_xlsx_single()
    make_xlsx_multisheet_linked()
    make_ppt_table_dup()
    print(f"\n총 파일: {len(list(OUT.iterdir()))}개 → {OUT.resolve()}")